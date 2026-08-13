from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(md_file_path, output_pptx_path):
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.readlines()

    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Agentes de IA Autónomos"
    subtitle.text = "El Futuro de la Colaboración Inteligente"

    # Process content
    current_slide = None
    
    for line in content:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith("# "):
            pass # Title covered
        elif line.startswith("## "):
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = current_slide.shapes.title
            title.text = line.replace("## ", "")
        elif line.startswith("### "):
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = current_slide.shapes.title
            title.text = line.replace("### ", "")
        elif current_slide:
            body = current_slide.shapes.placeholders[1]
            if line.startswith("* "):
                p = body.text_frame.add_paragraph()
                p.text = line.replace("* ", "")
                p.level = 1
            else:
                if not body.text:
                    body.text = line
                else:
                    body.text += "\n" + line

    prs.save(output_pptx_path)
    print(f"Presentación guardada en: {output_pptx_path}")

create_presentation(r'C:\Users\hp\Documents\PROYECTOS\JARVIS\presentacion_agentes_ia.md', r'C:\Users\hp\Documents\PROYECTOS\JARVIS\presentacion_agentes_ia.pptx')
