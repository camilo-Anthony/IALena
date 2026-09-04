from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Agentes de IA: El Futuro de la Autonomía"
    slide.placeholders[1].text = "JARVIS / IALena\nProyecto de Desarrollo"

    # Slide 1: What is an Agent
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "¿Qué es un Agente de IA?"
    slide.placeholders[1].text = (
        "Un sistema capaz de percibir su entorno, razonar, "
        "tomar decisiones y ejecutar acciones para lograr objetivos específicos."
    )

    # Slide 2: Key Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Componentes Clave"
    slide.placeholders[1].text = (
        "- Percepción (Sensores/Entradas)\n"
        "- Razonamiento (Modelo/LLM)\n"
        "- Memoria (Histórica/Durable)\n"
        "- Acción (Herramientas/APIs)"
    )

    # Slide 3: JARVIS Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "JARVIS: Enfoque Autónomo"
    slide.placeholders[1].text = (
        "- Serialización de tareas (TaskLedger)\n"
        "- Gestión de estado persistente\n"
        "- Interfaz dual (Visual/Técnica)\n"
        "- Autonomía de ejecución"
    )

    output_path = "C:/Users/hp/Documents/PROYECTOS/JARVIS/Presentacion_Agentes_IA.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    print(create_presentation())
